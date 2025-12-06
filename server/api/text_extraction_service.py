"""
Text Extraction Service
Extracts text from multiple file formats: PDF, DOCX, PPTX, XLSX, Images, TXT

Uses file_type_registry for centralized file type detection.
Supports Surya OCR for images with Tesseract fallback.
"""
import os
import io
import logging
from pathlib import Path
from typing import Optional, Tuple

# Core libraries
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook
import pytesseract
from PIL import Image
from langdetect import detect, LangDetectException
import chardet

# Import file type registry
from .file_type_registry import (
    FILE_TYPE_REGISTRY,
    FileCategory,
    get_extension,
    get_file_info,
    is_supported,
    requires_ocr,
    validate_file,
)

logger = logging.getLogger(__name__)


class TextExtractionError(Exception):
    """Custom exception for text extraction errors."""
    pass


class TextExtractor:
    """
    Unified text extraction from multiple file formats.
    
    Uses file_type_registry for file type detection and routing.
    Supported formats:
    - PDF (.pdf) - Native text extraction with OCR fallback
    - DOCX (.docx) - Microsoft Word documents
    - PPTX (.pptx) - Microsoft PowerPoint presentations
    - XLSX (.xlsx) - Microsoft Excel spreadsheets
    - Images (.png, .jpg, .jpeg, .gif, .bmp, .tiff, .webp) - OCR (Surya/Tesseract)
    - Text files (.txt, .md, .rtf, .csv, .json, .xml, .html)
    """
    
    # Use registry for extensions (backward compatibility)
    SUPPORTED_EXTENSIONS = {ext: info.category.value for ext, info in FILE_TYPE_REGISTRY.items()}
    
    # OCR language mapping (common languages)
    OCR_LANGUAGES = {
        'en': 'eng',
        'de': 'deu',
        'ar': 'ara',
        'fr': 'fra',
        'es': 'spa',
        'it': 'ita',
        'pt': 'por',
        'ru': 'rus',
        'zh': 'chi_sim',
        'ja': 'jpn',
        'ko': 'kor',
        'hi': 'hin',
        'tr': 'tur',
        'nl': 'nld',
        'pl': 'pol',
    }
    
    def __init__(self, ocr_lang: str = 'eng+deu+ara'):
        """
        Initialize the text extractor.
        
        Args:
            ocr_lang: Tesseract language codes for OCR (default: eng+deu+ara)
        """
        self.ocr_lang = ocr_lang
    
    def extract(self, file_content: bytes, filename: str, ocrspace_api_key: str = None) -> dict:
        """
        Extract text from a file.
        
        Args:
            file_content: Raw file bytes
            filename: Original filename with extension
            ocrspace_api_key: Optional OCR.space API key for AI-powered image OCR
            
        Returns:
            dict with keys:
                - text: Extracted text content
                - language: Detected language code
                - file_type: Detected file type
                - pages: Number of pages (if applicable)
                - word_count: Number of words extracted
                - metadata: Additional file metadata
        """
        ext = self._get_extension(filename)
        file_category = self.SUPPORTED_EXTENSIONS.get(ext)
        
        if not file_category:
            raise TextExtractionError(
                f"Unsupported file format: .{ext}. "
                f"Supported: {', '.join(self.SUPPORTED_EXTENSIONS.keys())}"
            )
        
        try:
            # Route based on extension for specific handlers
            if ext == 'pdf':
                text, pages, metadata = self._extract_pdf(file_content)
            elif ext == 'docx':
                text, pages, metadata = self._extract_docx(file_content)
            elif ext == 'doc':
                raise TextExtractionError(
                    "Legacy .doc format not supported. Please convert to .docx"
                )
            elif ext == 'pptx':
                text, pages, metadata = self._extract_pptx(file_content)
            elif ext == 'ppt':
                raise TextExtractionError(
                    "Legacy .ppt format not supported. Please convert to .pptx"
                )
            elif ext == 'xlsx':
                text, pages, metadata = self._extract_xlsx(file_content)
            elif ext == 'xls':
                raise TextExtractionError(
                    "Legacy .xls format not supported. Please convert to .xlsx"
                )
            elif file_category == 'image':
                # Pass OCR.space API key for AI-powered OCR
                text, pages, metadata = self._extract_image(file_content, ocrspace_api_key=ocrspace_api_key)
            elif file_category == 'text':
                text, pages, metadata = self._extract_text(file_content, ext)
            else:
                raise TextExtractionError(f"Unknown file type: {file_category}")


            
            # Detect language
            language = self._detect_language(text)
            
            # Count words
            word_count = len(text.split()) if text else 0
            
            return {
                'text': text.strip(),
                'language': language,
                'file_type': ext,
                'pages': pages,
                'word_count': word_count,
                'metadata': metadata,
            }
            
        except TextExtractionError:
            raise
        except Exception as e:
            logger.exception(f"Text extraction failed for {filename}")
            raise TextExtractionError(f"Failed to extract text: {str(e)}")
    
    def _get_extension(self, filename: str) -> str:
        """Get lowercase file extension without dot."""
        return Path(filename).suffix.lower().lstrip('.')
    
    def _detect_language(self, text: str) -> str:
        """Detect text language using langdetect."""
        if not text or len(text.strip()) < 20:
            return 'unknown'
        try:
            # Use first 1000 chars for faster detection
            sample = text[:1000]
            return detect(sample)
        except LangDetectException:
            return 'unknown'
    
    def _extract_pdf(self, content: bytes) -> Tuple[str, int, dict]:
        """Extract text from PDF using PyMuPDF, with OCR fallback."""
        doc = fitz.open(stream=content, filetype="pdf")
        
        text_parts = []
        pages = len(doc)
        metadata = {
            'title': doc.metadata.get('title', ''),
            'author': doc.metadata.get('author', ''),
            'subject': doc.metadata.get('subject', ''),
        }
        
        for page_num, page in enumerate(doc):
            # Try native text extraction first
            page_text = page.get_text("text")
            
            # If no text found, try OCR
            if not page_text.strip():
                try:
                    # Render page as image for OCR
                    pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    page_text = pytesseract.image_to_string(img, lang=self.ocr_lang)
                except Exception as e:
                    logger.warning(f"OCR failed on page {page_num + 1}: {e}")
            
            if page_text.strip():
                text_parts.append(f"--- Page {page_num + 1} ---\n{page_text}")
        
        doc.close()
        return "\n\n".join(text_parts), pages, metadata
    
    def _extract_docx(self, content: bytes) -> Tuple[str, int, dict]:
        """Extract text from DOCX document."""
        doc = DocxDocument(io.BytesIO(content))
        
        text_parts = []
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Extract tables
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                table_text.append(" | ".join(row_text))
            if table_text:
                text_parts.append("\n[Table]\n" + "\n".join(table_text))
        
        # Get metadata
        core_props = doc.core_properties
        metadata = {
            'title': core_props.title or '',
            'author': core_props.author or '',
            'subject': core_props.subject or '',
        }
        
        return "\n\n".join(text_parts), 1, metadata
    
    def _extract_pptx(self, content: bytes) -> Tuple[str, int, dict]:
        """Extract text from PowerPoint presentation."""
        prs = Presentation(io.BytesIO(content))
        
        text_parts = []
        slides = len(prs.slides)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # Handle tables
                if shape.has_table:
                    table_text = []
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        table_text.append(" | ".join(row_text))
                    if table_text:
                        slide_text.append("[Table]\n" + "\n".join(table_text))
            
            if slide_text:
                text_parts.append(f"--- Slide {slide_num} ---\n" + "\n".join(slide_text))
        
        metadata = {
            'title': prs.core_properties.title or '',
            'author': prs.core_properties.author or '',
        }
        
        return "\n\n".join(text_parts), slides, metadata
    
    def _extract_xlsx(self, content: bytes) -> Tuple[str, int, dict]:
        """Extract text from Excel spreadsheet."""
        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        
        text_parts = []
        sheets = len(wb.sheetnames)
        
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_text = [f"--- Sheet: {sheet_name} ---"]
            
            for row in sheet.iter_rows():
                row_values = []
                for cell in row:
                    if cell.value is not None:
                        row_values.append(str(cell.value))
                if row_values:
                    sheet_text.append(" | ".join(row_values))
            
            if len(sheet_text) > 1:  # Has content beyond header
                text_parts.append("\n".join(sheet_text))
        
        wb.close()
        
        metadata = {'sheets': wb.sheetnames}
        return "\n\n".join(text_parts), sheets, metadata
    
    def _extract_image(self, content: bytes, use_ai_ocr: bool = True, ocrspace_api_key: str = None) -> Tuple[str, int, dict]:
        """Extract text from image using OCR.
        
        Args:
            content: Image bytes
            use_ai_ocr: If True, try OCR.space API first (default: True)
            ocrspace_api_key: OCR.space API key for AI OCR
        """
        # Try OCR.space API first if API key is provided
        if use_ai_ocr and ocrspace_api_key:
            try:
                from .ocrspace_service import extract_text_with_ocrspace
                
                logger.info("Using OCR.space API for image extraction")
                text, metadata = extract_text_with_ocrspace(content, ocrspace_api_key)
                
                # Get image dimensions for metadata
                img = Image.open(io.BytesIO(content))
                metadata['width'] = img.width
                metadata['height'] = img.height
                metadata['format'] = img.format
                
                return text, 1, metadata
            except Exception as e:
                logger.warning(f"OCR.space failed, falling back to Tesseract: {e}")
        
        # Fallback to Tesseract
        logger.info("Using Tesseract OCR for image extraction")
        img = Image.open(io.BytesIO(content))
        
        # Convert to RGB if necessary
        if img.mode in ('RGBA', 'P'):
            img = img.convert('RGB')
        
        try:
            text = pytesseract.image_to_string(img, lang=self.ocr_lang)
        except Exception as e:
            error_msg = str(e)
            if 'tesseract is not installed' in error_msg.lower() or 'not in your path' in error_msg.lower():
                # Provide helpful error message
                if not ocrspace_api_key:
                    raise TextExtractionError(
                        "Image OCR requires either: (1) Add your OCR.space API key in Settings → API Keys, or "
                        "(2) Install Tesseract OCR on your system. Get a free OCR.space key at: https://ocr.space/ocrapi"
                    )
            raise TextExtractionError(f"OCR failed: {error_msg}")
        
        metadata = {
            'method': 'tesseract',
            'width': img.width,
            'height': img.height,
            'format': img.format,
        }
        
        return text, 1, metadata


    
    def _extract_text(self, content: bytes, ext: str) -> Tuple[str, int, dict]:
        """Extract text from plain text files."""
        # Detect encoding
        detected = chardet.detect(content)
        encoding = detected.get('encoding', 'utf-8') or 'utf-8'
        
        try:
            text = content.decode(encoding)
        except (UnicodeDecodeError, LookupError):
            # Fallback to utf-8 with error handling
            text = content.decode('utf-8', errors='replace')
        
        metadata = {
            'encoding': encoding,
            'format': ext,
        }
        
        return text, 1, metadata


# Singleton instance for easy import
extractor = TextExtractor()


def extract_text_from_file(file_content: bytes, filename: str, ocrspace_api_key: str = None) -> dict:
    """
    Convenience function to extract text from a file.
    
    Args:
        file_content: Raw file bytes
        filename: Original filename with extension
        ocrspace_api_key: Optional OCR.space API key for AI-powered image OCR
        
    Returns:
        dict with text, language, file_type, pages, word_count, metadata
    """
    return extractor.extract(file_content, filename, ocrspace_api_key=ocrspace_api_key)

